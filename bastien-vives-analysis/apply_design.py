"""
apply_design.py

Reads 04_presentation_draft.pptx and applies a consistent design system,
saving the result as 05_presentation_styled.pptx.

Design Rules
------------
Primary accent  : #e63946  (red)
Secondary / nav : #1a1a2e  (deep navy)
Background      : #FAFAF8  (cream white)
Text            : #1a1a2e
Font            : 나눔고딕 → Malgun Gothic → Arial (fallback chain)
Title slide (1) : navy background, white text
Section dividers: red background, white text  (slides whose title starts with "[Unit")
Content slides  : cream-white background, navy text
Table header    : navy background, white text
Slide numbers   : bottom-right
Margins         : 0.5 in consistent
"""

import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ────────────────────────────────────────────────────────────────────
BASE = os.path.expanduser("~/research-output/바스티앙_비베스_작품분석")
INPUT  = os.path.join(BASE, "04_presentation_draft.pptx")
OUTPUT = os.path.join(BASE, "05_presentation_styled.pptx")

# ── Palette ──────────────────────────────────────────────────────────────────
C_RED    = RGBColor(0xe6, 0x39, 0x46)   # #e63946
C_NAVY   = RGBColor(0x1a, 0x1a, 0x2e)   # #1a1a2e
C_CREAM  = RGBColor(0xFA, 0xFA, 0xF8)   # #FAFAF8
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED  = RGBColor(0xAA, 0xAA, 0xBB)   # slide-number muted text

# ── Font preference list ──────────────────────────────────────────────────────
FONT_PREFERENCE = ["나눔고딕", "Malgun Gothic", "Arial"]


# ─────────────────────────────────────────────────────────────────────────────
# Helper: set background solid fill
# ─────────────────────────────────────────────────────────────────────────────
def set_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Helper: apply font to a run element, with graceful fallback
# ─────────────────────────────────────────────────────────────────────────────
def apply_font(run, size_pt=None, bold=None, color: RGBColor = None,
               font_name: str = FONT_PREFERENCE[0]):
    try:
        run.font.name = font_name
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    except Exception:
        pass  # graceful fallback — keep whatever was there


# ─────────────────────────────────────────────────────────────────────────────
# Helper: apply font to every run in a text frame
# ─────────────────────────────────────────────────────────────────────────────
def restyle_text_frame(tf, color: RGBColor = None, font_name: str = FONT_PREFERENCE[0]):
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                run.font.name = font_name
                if color is not None:
                    run.font.color.rgb = color
            except Exception:
                pass


# ─────────────────────────────────────────────────────────────────────────────
# Helper: detect slide role
# ─────────────────────────────────────────────────────────────────────────────
def get_slide_role(slide_index: int, slide) -> str:
    """
    Returns one of: 'title', 'section', 'content'
    - slide index 0  → title
    - slides whose title text starts with '[Unit' → section divider
    - everything else → content
    """
    if slide_index == 0:
        return "title"

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.startswith("[Unit"):
                return "section"

    return "content"


# ─────────────────────────────────────────────────────────────────────────────
# Helper: find or create slide-number text box
# We identify it by looking for a small text-box at bottom-right containing
# only a number.  If not found we don't add a new one (preserve originals).
# ─────────────────────────────────────────────────────────────────────────────
def update_slide_number_style(slide, prs_width: Emu, prs_height: Emu,
                               text_color: RGBColor):
    """
    Re-style an existing slide-number box if it exists.
    The original build_pptx.py adds a TextBox at (W-1in, H-0.4in).
    We locate it by position proximity and short numeric text.
    """
    threshold = Inches(1.5)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text.isdigit():
            continue
        # Check position: bottom-right quadrant
        if shape.left > prs_width - threshold and shape.top > prs_height - threshold:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        run.font.color.rgb = text_color
                        run.font.name = FONT_PREFERENCE[0]
                    except Exception:
                        pass
            return


# ─────────────────────────────────────────────────────────────────────────────
# Per-role styling functions
# ─────────────────────────────────────────────────────────────────────────────

def style_title_slide(slide, prs_width, prs_height):
    """Slide 1: deep navy bg, white text, red accent line."""
    set_background(slide, C_NAVY)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            # Re-color accent shapes (the red bar should stay red)
            try:
                fill = shape.fill
                if fill.type is not None:
                    rgb = fill.fore_color.rgb
                    # Keep red accent bar as-is; re-color other shapes to navy
                    if rgb != C_RED:
                        fill.solid()
                        fill.fore_color.rgb = C_NAVY
            except Exception:
                pass
            continue

        text = shape.text_frame.text.strip()
        # Slide-number box
        if text.isdigit():
            update_slide_number_style(slide, prs_width, prs_height, C_MUTED)
            continue

        # All other text → white
        restyle_text_frame(shape.text_frame, color=C_WHITE,
                           font_name=FONT_PREFERENCE[0])


def style_section_slide(slide, prs_width, prs_height):
    """Section divider: red background, white text."""
    set_background(slide, C_RED)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            try:
                fill = shape.fill
                if fill.type is not None:
                    fill.solid()
                    fill.fore_color.rgb = C_RED
                shape.line.fill.background()
            except Exception:
                pass
            continue

        text = shape.text_frame.text.strip()
        if text.isdigit():
            update_slide_number_style(slide, prs_width, prs_height, C_WHITE)
            continue

        restyle_text_frame(shape.text_frame, color=C_WHITE,
                           font_name=FONT_PREFERENCE[0])

        # Increase font size for the main title text in section slides
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text.strip()
                if txt.startswith("[Unit") and len(txt) > 5:
                    try:
                        current_size = run.font.size
                        if current_size is None or current_size < Pt(28):
                            run.font.size = Pt(30)
                        run.font.bold = True
                    except Exception:
                        pass


def style_content_slide(slide, prs_width, prs_height):
    """Standard content slide: cream-white bg, navy text, red accents."""
    set_background(slide, C_CREAM)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            try:
                fill = shape.fill
                if fill.type is None:
                    continue
                rgb = fill.fore_color.rgb
                # Keep red accents red
                if rgb == C_RED or rgb == RGBColor(0xe6, 0x39, 0x46):
                    continue
                # Remap title-bar colors to navy or light cream
                # Original title bar color: 0x1a1a2e (navy) → keep
                # Original light bar: 0xF5F5F7 or 0xF0F4FF → remap to cream
                if rgb in (RGBColor(0xF5, 0xF5, 0xF7),
                           RGBColor(0xF0, 0xF4, 0xFF),
                           RGBColor(0xF8, 0xF9, 0xFF),
                           RGBColor(0xFF, 0xFF, 0xFF)):
                    fill.solid()
                    fill.fore_color.rgb = C_CREAM
                # navy stays navy
            except Exception:
                pass
            continue

        text = shape.text_frame.text.strip()
        if text.isdigit():
            update_slide_number_style(slide, prs_width, prs_height, C_MUTED)
            continue

        # Detect if this is a title text box (large font, near top)
        is_title_area = shape.top is not None and shape.top < Inches(1.5)

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    run.font.name = FONT_PREFERENCE[0]
                    current_color = None
                    try:
                        current_color = run.font.color.rgb
                    except Exception:
                        pass

                    # White text in navy bars should stay white
                    if current_color == C_WHITE:
                        continue
                    # Grey/muted → keep muted for unit tags
                    if current_color in (RGBColor(0x99, 0x99, 0x99),
                                         RGBColor(0xCC, 0xCC, 0xCC)):
                        run.font.color.rgb = C_MUTED
                        continue

                    # Red accent text → keep red
                    if current_color == C_RED:
                        continue

                    # Everything else → navy
                    run.font.color.rgb = C_NAVY
                except Exception:
                    pass


# ─────────────────────────────────────────────────────────────────────────────
# Table header styling (for table shapes, if any)
# ─────────────────────────────────────────────────────────────────────────────
def style_tables(slide):
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        # First row = header
        for cell in table.rows[0].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_NAVY
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        run.font.color.rgb = C_WHITE
                        run.font.bold = True
                        run.font.name = FONT_PREFERENCE[0]
                    except Exception:
                        pass
        # Other rows
        for row in table.rows[1:]:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            run.font.color.rgb = C_NAVY
                            run.font.name = FONT_PREFERENCE[0]
                        except Exception:
                            pass


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def apply_design(input_path: str, output_path: str):
    prs = Presentation(input_path)
    W = prs.slide_width
    H = prs.slide_height

    for idx, slide in enumerate(prs.slides):
        role = get_slide_role(idx, slide)

        if role == "title":
            style_title_slide(slide, W, H)
        elif role == "section":
            style_section_slide(slide, W, H)
        else:
            style_content_slide(slide, W, H)

        # Apply table styling regardless of role
        style_tables(slide)

    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"Saved: {output_path}")
    print(f"Size : {size_kb:,} KB ({os.path.getsize(output_path):,} bytes)")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    apply_design(INPUT, OUTPUT)
