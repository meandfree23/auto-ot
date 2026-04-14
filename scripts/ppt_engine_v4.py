#!/usr/bin/env python3
"""
OT 분석 리포트 HTML 생성 엔진 v5.0
World-class design: refined glassmorphism · Chart.js · phase timeline · micro-animations
"""
import os, sys, re, json, argparse
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.append(str(ROOT))

# ─── Category metadata ────────────────────────────────────────────────────────
CATEGORY_META = {
    "Market Intelligence": {
        "emoji": "📊", "icon": "bar-chart-2",
        "image": "1551288049-bebda4e38f71",
        "label_ko": "시장 분석"
    },
    "Strategic Logic": {
        "emoji": "🎯", "icon": "target",
        "image": "1552664730-d307ca884978",
        "label_ko": "전략 논리"
    },
    "Creative & Visual": {
        "emoji": "🎨", "icon": "palette",
        "image": "1542744173-8e7e53415bb0",
        "label_ko": "크리에이티브"
    },
    "Gap Check & Risks": {
        "emoji": "⚠️", "icon": "alert-triangle",
        "image": "1507003211169-0a1dd7228f2d",
        "label_ko": "리스크 점검"
    },
}

# ─── Dynamic theme ────────────────────────────────────────────────────────────
def get_dynamic_theme(title):
    themes = [
        {"hex": "#3b82f6", "rgb": "59,130,246"},   # Blue
        {"hex": "#8b5cf6", "rgb": "139,92,246"},   # Purple
        {"hex": "#f97316", "rgb": "249,115,22"},   # Orange
        {"hex": "#10b981", "rgb": "16,185,129"},   # Green
        {"hex": "#ec4899", "rgb": "236,72,153"},   # Pink
        {"hex": "#f43f5e", "rgb": "244,63,94"},    # Red
        {"hex": "#06b6d4", "rgb": "6,182,212"},    # Cyan
    ]
    tl = title.lower()
    if any(k in tl for k in ["테크","it","플랫폼","디지털","앱"]): return themes[0]
    if any(k in tl for k in ["전략","b2b","금융","컨설"]): return themes[1]
    if any(k in tl for k in ["커머스","명품","프리미엄","쇼핑"]): return themes[2]
    if any(k in tl for k in ["친환경","esg","식품","반려","펫","pet"]): return themes[3]
    if any(k in tl for k in ["뷰티","화장품","여성","코스메틱"]): return themes[4]
    if any(k in tl for k in ["게임","메타버스","엔터","web3"]): return themes[5]
    if any(k in tl for k in ["헬스","의료","건강","웰니스"]): return themes[6]
    idx = sum(ord(c) for c in title) % len(themes)
    return themes[idx]

# ─── Score computation for radar chart ───────────────────────────────────────
def compute_radar_scores(categories):
    order = ["Market Intelligence","Strategic Logic","Creative & Visual","Gap Check & Risks"]
    scores = []
    for cat in order:
        secs = categories.get(cat, [])
        if not secs:
            scores.append(25); continue
        s = 0
        for sec in secs:
            s += 20
            for line in sec.get("contents", []):
                if line.startswith(("- ","* ")): s += 6
                elif line.startswith("|") and "---" not in line: s += 10
                elif line.startswith("> "): s += 8
                elif line.startswith("!["): s += 7
        scores.append(min(98, s))
    return scores

# ─── Markdown inline rendering ────────────────────────────────────────────────
def md_inline(text):
    """Convert **bold**, *italic* to HTML."""
    text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', text)
    return text

# ─── Table rendering ──────────────────────────────────────────────────────────
def is_phase_table(rows):
    for r in rows:
        if re.search(r'(단계|Phase\s*\d)', r, re.I): return True
    return False

def render_phase_timeline(rows):
    steps = []
    for r in rows:
        if "---" in r: continue
        cells = [c.strip() for c in r.split("|") if c.strip()]
        if not cells or len(cells) < 2: continue
        # Skip header row (column label names, not actual phase data)
        if cells[0].lower() in ('phase', '단계', '구분', 'category', 'phase / 단계'):
            continue
        phase = md_inline(cells[0])
        action = md_inline(cells[1]) if len(cells) > 1 else ""
        kpi = md_inline(cells[2]) if len(cells) > 2 else ""
        steps.append((phase, action, kpi))
    if not steps: return ""
    html = '<div class="phase-timeline">'
    for i, (phase, action, kpi) in enumerate(steps):
        html += f'''
<div class="phase-step">
  <div class="phase-num">{i+1}</div>
  <div class="phase-body">
    <div class="phase-title">{phase}</div>
    <div class="phase-action">{action}</div>
    {f'<div class="phase-kpi">🎯 {kpi}</div>' if kpi else ""}
  </div>
</div>'''
    html += '</div>'
    return html

def render_table(rows):
    """Render markdown table rows as styled HTML."""
    header_done = False
    html = '<div class="tbl-wrap"><table class="md-table">'
    for r in rows:
        if "---" in r: header_done = True; continue
        cells = [c.strip() for c in r.split("|") if c.strip()]
        if not cells: continue
        tag = "th" if not header_done else "td"
        tds = "".join(f'<{tag}>{md_inline(c)}</{tag}>' for c in cells)
        html += f"<tr>{tds}</tr>"
        if not header_done: header_done = True
    html += "</table></div>"
    return html

# ─── Content block rendering ──────────────────────────────────────────────────
def render_content_blocks(lines, chart_id_ref):
    html = ""
    i = 0
    while i < len(lines):
        line = lines[i]

        # Image
        if line.startswith("!["):
            m_url = re.search(r'\((.+?)\)', line)
            m_alt = re.search(r'\[(.+?)\]', line)
            url = m_url.group(1) if m_url else ""
            alt = m_alt.group(1) if m_alt else ""
            if url:
                html += f'<div class="img-block"><img src="{url}" alt="{alt}" loading="lazy"><div class="img-label">{alt}</div></div>'
            i += 1; continue

        # Blockquote
        if line.startswith("> "):
            html += f'<blockquote class="ot-quote"><p>{md_inline(line[2:])}</p></blockquote>'
            i += 1; continue

        # Table — collect all consecutive table rows
        if line.startswith("|"):
            table_rows = []
            while i < len(lines) and lines[i].startswith("|"):
                table_rows.append(lines[i])
                i += 1
            if is_phase_table(table_rows):
                html += render_phase_timeline(table_rows)
            else:
                html += render_table(table_rows)
            continue

        # Bullet
        if line.startswith(("- ", "* ")):
            # collect consecutive bullets
            items = []
            while i < len(lines) and lines[i].startswith(("- ", "* ")):
                items.append(md_inline(lines[i][2:]))
                i += 1
            html += '<ul class="ot-bullets">' + "".join(f'<li>{it}</li>' for it in items) + '</ul>'
            continue

        # Skip horizontal dividers
        if line.strip() in ('---', '- - -', '***', '* * *', '___'):
            i += 1; continue
        # Plain text
        if line.strip():
            html += f'<p class="ot-para">{md_inline(line)}</p>'
        i += 1
    return html

# ─── Card rendering ───────────────────────────────────────────────────────────
_chart_counter = [0]

def render_card(sec, cat_name):
    is_risk = "Gap" in cat_name or "Risk" in cat_name
    tag_cls = "tag-risk" if is_risk else "tag-insight"
    tag_txt = "⚠️ RISK" if is_risk else "💡 INSIGHT"
    title = md_inline(sec.get("title", ""))
    contents = sec.get("contents", [])
    _chart_counter[0] += 1
    body = render_content_blocks(contents, _chart_counter[0])
    return f'''
<div class="ot-card">
  <div class="card-header">
    <span class="card-tag {tag_cls}">{tag_txt}</span>
  </div>
  <h3 class="card-title">{title}</h3>
  <div class="card-body">{body}</div>
</div>'''

# ─── Category section rendering ───────────────────────────────────────────────
def render_category(cat_name, sections):
    meta = CATEGORY_META.get(cat_name, {"emoji": "📌", "image": "1551288049-bebda4e38f71", "label_ko": cat_name})
    img_url = f"https://images.unsplash.com/photo-{meta['image']}?w=1600&q=70&fm=webp&fit=crop&auto=format"
    cards_html = "\n".join(render_card(s, cat_name) for s in sections)
    return f'''
<section class="cat-section" id="cat-{cat_name.replace(' ','-').lower()}">
  <div class="cat-banner" style="background-image:url('{img_url}')">
    <div class="cat-banner-overlay"></div>
    <div class="cat-banner-content">
      <span class="cat-emoji">{meta["emoji"]}</span>
      <h2 class="cat-title">{cat_name}</h2>
    </div>
  </div>
  <div class="cards-grid">
    {cards_html}
  </div>
</section>'''

# ─── HTML Template ────────────────────────────────────────────────────────────
HTML_TEMPLATE = r"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>[[TITLE]]</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap" rel="stylesheet">
<script src="https://cdn.tailwindcss.com"></script>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<script src="https://unpkg.com/lucide@latest/dist/umd/lucide.min.js"></script>
<style>
:root {
  --accent: [[BRAND_HEX]];
  --accent-rgb: [[BRAND_RGB]];
  --bg: #030303;
  --surface: rgba(255,255,255,0.035);
  --surface-hov: rgba(255,255,255,0.06);
  --border: rgba(255,255,255,0.07);
  --border-hov: rgba([[BRAND_RGB]],0.3);
  --text-1: #f5f5f7;
  --text-2: #a1a1a6;
  --text-3: #6e6e73;
  --r-card: 20px;
  --glow: 0 20px 80px -20px rgba([[BRAND_RGB]],.35), 0 0 0 1px rgba([[BRAND_RGB]],.08);
}
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
html { scroll-behavior: smooth; }
body {
  font-family: 'Inter','Noto Sans KR',sans-serif;
  background: var(--bg);
  color: var(--text-1);
  -webkit-font-smoothing: antialiased;
  background-image:
    radial-gradient(ellipse 90% 60% at 10% -10%, rgba([[BRAND_RGB]],.13) 0%, transparent 60%),
    radial-gradient(ellipse 70% 50% at 90% 110%, rgba([[BRAND_RGB]],.07) 0%, transparent 60%);
  background-attachment: fixed;
}
/* ── Nav ────────────────────────────── */
.site-nav {
  position: fixed; top: 0; left: 0; right: 0; z-index: 100;
  display: flex; align-items: center; justify-content: space-between;
  padding: 0 2rem; height: 60px;
  background: rgba(3,3,3,.7);
  backdrop-filter: blur(24px);
  border-bottom: 1px solid var(--border);
}
.nav-brand { display: flex; align-items: center; gap: .6rem; font-size: .85rem; font-weight: 700; letter-spacing: .08em; }
.nav-dot { width: 8px; height: 8px; border-radius: 50%; background: var(--accent); box-shadow: 0 0 10px rgba([[BRAND_RGB]],.8); }
.nav-meta { font-size: .72rem; color: var(--text-3); font-family: monospace; letter-spacing: .05em; }
/* ── Hero ───────────────────────────── */
.hero {
  min-height: 92vh; padding: 120px 5vw 80px;
  display: grid; grid-template-columns: 1fr 380px; gap: 4rem; align-items: center;
}
@media(max-width:900px){ .hero { grid-template-columns: 1fr; padding: 100px 5vw 60px; } }
.hero-badge {
  display: inline-block; margin-bottom: 1.5rem;
  padding: .3rem 1rem; border-radius: 999px;
  border: 1px solid rgba([[BRAND_RGB]],.3);
  background: rgba([[BRAND_RGB]],.08);
  font-size: .72rem; font-weight: 700; letter-spacing: .15em; text-transform: uppercase;
  color: var(--accent);
}
.hero-title {
  font-size: clamp(2.4rem,5vw,4.5rem);
  font-weight: 900; line-height: 1.08; letter-spacing: -.04em;
  background: linear-gradient(160deg, #fff 30%, rgba(255,255,255,.45));
  -webkit-background-clip: text; -webkit-text-fill-color: transparent;
  margin-bottom: 1.5rem;
}
.hero-sub { font-size: 1.05rem; color: var(--text-2); line-height: 1.7; max-width: 560px; margin-bottom: 2.5rem; }
.hero-chips { display: flex; flex-wrap: wrap; gap: .6rem; }
.chip {
  padding: .35rem .85rem; border-radius: 8px;
  background: rgba(255,255,255,.05); border: 1px solid var(--border);
  font-size: .75rem; color: var(--text-2); font-weight: 500; display: flex; align-items: center; gap: .4rem;
}
.chip-accent { color: var(--accent); border-color: rgba([[BRAND_RGB]],.2); background: rgba([[BRAND_RGB]],.06); }
/* ── Radar chart container ──────────── */
.radar-wrap {
  display: flex; flex-direction: column; align-items: center; justify-content: center;
  padding: 2rem; border-radius: 28px;
  background: var(--surface); border: 1px solid var(--border);
  backdrop-filter: blur(40px);
}
.radar-wrap canvas { max-width: 300px; max-height: 300px; }
.radar-label { font-size: .72rem; color: var(--text-3); margin-top: 1rem; text-align: center; letter-spacing: .1em; text-transform: uppercase; }
/* ── Section banner ─────────────────── */
.cat-section { margin-bottom: 6rem; padding: 0 5vw; }
.cat-banner {
  width: 100%; height: 220px; border-radius: 24px; overflow: hidden;
  position: relative; margin-bottom: 3rem;
  background-size: cover; background-position: center;
}
.cat-banner-overlay {
  position: absolute; inset: 0;
  background: linear-gradient(120deg, rgba(3,3,3,.85) 0%, rgba(3,3,3,.4) 100%);
}
.cat-banner-content {
  position: absolute; bottom: 2rem; left: 2.5rem;
  display: flex; align-items: center; gap: 1rem;
}
.cat-emoji { font-size: 2.5rem; line-height: 1; }
.cat-title {
  font-size: 2rem; font-weight: 800; letter-spacing: -.03em;
  background: linear-gradient(90deg, #fff, rgba(255,255,255,.7));
  -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
/* ── Cards grid ─────────────────────── */
.cards-grid {
  display: flex; flex-direction: column; gap: 2rem;
}
.ot-card {
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: var(--r-card);
  padding: 2.5rem;
  backdrop-filter: blur(30px);
  transition: transform .4s cubic-bezier(.16,1,.3,1), border-color .3s, box-shadow .3s;
  animation: cardIn .7s cubic-bezier(.16,1,.3,1) both;
}
@keyframes cardIn {
  from { opacity: 0; transform: translateY(24px); filter: blur(6px); }
  to   { opacity: 1; transform: none; filter: none; }
}
.ot-card:hover {
  border-color: var(--border-hov);
  box-shadow: var(--glow);
  transform: translateY(-6px);
}
.card-header { display: flex; align-items: center; gap: .6rem; margin-bottom: 1.5rem; }
.card-tag {
  font-size: .68rem; font-weight: 700; letter-spacing: .12em; text-transform: uppercase;
  padding: .25rem .7rem; border-radius: 6px;
}
.tag-insight { background: rgba([[BRAND_RGB]],.12); color: var(--accent); border: 1px solid rgba([[BRAND_RGB]],.25); }
.tag-risk { background: rgba(244,63,94,.1); color: #f87171; border: 1px solid rgba(244,63,94,.25); }
.card-title { font-size: 1rem; font-weight: 600; line-height: 1.35; letter-spacing: -.015em; margin-bottom: 1.4rem; color: var(--text-1); display: none; }
.card-body { display: flex; flex-direction: column; gap: 1.4rem; }
/* ── Global text quality ────────────── */
p, li, td, blockquote {
  word-break: keep-all; overflow-wrap: break-word;
  text-rendering: optimizeLegibility;
  font-feature-settings: 'kern' 1, 'liga' 1;
}
strong { color: var(--text-1); font-weight: 700; }

/* ── Blockquote ─────────────────────── */
.ot-quote {
  border-left: 2px solid var(--accent);
  padding: .9rem 1.25rem;
  border-radius: 0 10px 10px 0;
  background: rgba([[BRAND_RGB]],.05);
}
.ot-quote p {
  font-size: .93rem;
  color: var(--text-1); line-height: 1.8;
  letter-spacing: -.015em; font-style: italic;
}

/* ── Bullets ────────────────────────── */
.ot-bullets {
  list-style: none; padding: 0;
  display: flex; flex-direction: column; gap: 0;
  border: 1px solid var(--border); border-radius: 14px; overflow: hidden;
}
.ot-bullets li {
  position: relative;
  padding: .85rem 1rem .85rem 1.5rem;
  font-size: .875rem; color: var(--text-2); line-height: 1.75;
  letter-spacing: -.01em; word-break: keep-all;
  border-bottom: 1px solid var(--border);
  transition: background .2s;
}
.ot-bullets li:last-child { border-bottom: none; }
.ot-bullets li:hover { background: rgba(255,255,255,.025); }
.ot-bullets li::before {
  content: '';
  position: absolute; left: .7rem; top: 1.15rem;
  width: 5px; height: 5px; border-radius: 50%;
  background: var(--accent);
  box-shadow: 0 0 6px rgba([[BRAND_RGB]],.6);
}
/* strong inside bullets: accent label style */
.ot-bullets li strong {
  display: inline-block;
  color: var(--text-1); font-weight: 700;
  background: rgba([[BRAND_RGB]],.1);
  border: 1px solid rgba([[BRAND_RGB]],.2);
  padding: .02em .45em; border-radius: 5px;
  font-size: .84em; letter-spacing: 0;
  margin-right: .15em;
}
/* First bullet: lead emphasis */
.ot-bullets li:first-child {
  color: var(--text-1); font-weight: 500; font-size: .9rem;
  background: rgba([[BRAND_RGB]],.04);
}

/* ── Table ──────────────────────────── */
.tbl-wrap {
  overflow-x: auto; border-radius: 16px;
  border: 1px solid var(--border); margin: .75rem 0;
  box-shadow: 0 4px 24px rgba(0,0,0,.2);
}
.md-table { width: 100%; border-collapse: collapse; font-size: .84rem; min-width: 480px; }
.md-table th {
  padding: .75rem 1rem; text-align: left; font-weight: 700; font-size: .71rem;
  letter-spacing: .08em; text-transform: uppercase; color: var(--accent);
  background: rgba([[BRAND_RGB]],.09); border-bottom: 1px solid rgba([[BRAND_RGB]],.18);
  word-break: keep-all; white-space: nowrap;
}
.md-table td {
  padding: .9rem 1rem; color: var(--text-2); word-break: keep-all;
  border-bottom: 1px solid rgba(255,255,255,.04); line-height: 1.7; vertical-align: top;
  font-size: .83rem;
}
.md-table td:first-child {
  color: var(--text-1); font-weight: 600; font-size: .8rem;
  white-space: nowrap; max-width: 140px;
}
.md-table td strong { color: var(--accent); font-weight: 700; background: none; padding: 0; }
.md-table tr:last-child td { border-bottom: none; }
.md-table tr:nth-child(even) td { background: rgba(255,255,255,.015); }
.md-table tr:hover td { background: rgba([[BRAND_RGB]],.04); }

/* ── Phase timeline ─────────────────── */
.phase-timeline { display: flex; flex-direction: column; gap: 0; margin: .75rem 0; }
.phase-step {
  display: flex; gap: 1.1rem; position: relative; padding-bottom: 1.6rem;
}
.phase-step:last-child { padding-bottom: 0; }
.phase-step::before {
  content: ''; position: absolute; left: 15px; top: 34px; bottom: 0; width: 1px;
  background: linear-gradient(to bottom, rgba([[BRAND_RGB]],.5), transparent);
}
.phase-step:last-child::before { display: none; }
.phase-num {
  flex-shrink: 0; width: 32px; height: 32px; border-radius: 50%;
  background: linear-gradient(135deg, var(--accent), rgba([[BRAND_RGB]],.7));
  color: #fff; font-size: .82rem; font-weight: 800;
  display: flex; align-items: center; justify-content: center;
  box-shadow: 0 0 20px rgba([[BRAND_RGB]],.5); z-index: 1;
}
.phase-body { flex: 1; padding-top: .2rem; }
.phase-title {
  font-size: .9rem; font-weight: 700; color: var(--text-1);
  margin-bottom: .4rem; line-height: 1.4; word-break: keep-all;
}
.phase-title strong { color: var(--accent); background: none; padding: 0; }
.phase-action {
  font-size: .84rem; color: var(--text-2); line-height: 1.7;
  margin-bottom: .4rem; word-break: keep-all;
}
.phase-kpi {
  font-size: .76rem; color: var(--accent);
  background: rgba([[BRAND_RGB]],.07); border: 1px solid rgba([[BRAND_RGB]],.18);
  border-radius: 6px; padding: .3rem .8rem; display: inline-block; margin-top: .3rem;
  word-break: keep-all;
}

/* ── Image block ────────────────────── */
.img-block { border-radius: 16px; overflow: hidden; margin: .75rem 0; position: relative; }
.img-block img {
  width: 100%; max-height: 240px; object-fit: cover; display: block;
  transition: transform .7s cubic-bezier(.16,1,.3,1);
}
.img-block:hover img { transform: scale(1.04); }
.img-label {
  position: absolute; bottom: 0; left: 0; right: 0;
  background: linear-gradient(to top, rgba(0,0,0,.75), transparent);
  padding: 1.8rem .9rem .6rem; font-size: .72rem; color: rgba(255,255,255,.65);
  letter-spacing: .04em;
}

/* ── Paragraph ──────────────────────── */
.ot-para {
  font-size: .9rem; color: var(--text-2); line-height: 1.75;
  letter-spacing: -.01em; word-break: keep-all;
}
.ot-para strong { color: var(--accent); background: rgba([[BRAND_RGB]],.08); padding: .05em .3em; border-radius: 4px; }

/* ── Divider ─────────────────────────── */
.ot-divider { border: none; border-top: 1px solid var(--border); margin: 1rem 0; }
/* ── Footer ─────────────────────────── */
.site-footer {
  margin-top: 6rem; padding: 3rem 5vw;
  border-top: 1px solid var(--border);
  display: flex; flex-wrap: wrap; gap: 1rem; align-items: center; justify-content: space-between;
}
.footer-left { font-size: .75rem; color: var(--text-3); line-height: 1.8; }
.footer-right { display: flex; gap: .5rem; }
.footer-chip {
  padding: .3rem .8rem; border-radius: 8px;
  background: var(--surface); border: 1px solid var(--border);
  font-size: .72rem; color: var(--text-3); font-family: monospace;
}
/* ── Scroll progress ────────────────── */
#scroll-progress {
  position: fixed; top: 60px; left: 0; height: 2px; width: 0%;
  background: linear-gradient(90deg, var(--accent), rgba([[BRAND_RGB]],.3));
  z-index: 200; transition: width .1s;
}
</style>
</head>
<body>
<!-- Scroll progress bar -->
<div id="scroll-progress"></div>

<!-- Nav -->
<nav class="site-nav">
  <div class="nav-brand">
    <div class="nav-dot"></div>
    <span>Antigravity Intelligence</span>
  </div>
  <span class="nav-meta">[[JOB_ID]]</span>
</nav>

<!-- Hero -->
<section class="hero">
  <div class="hero-left">
    <div class="hero-badge">[[BADGE_LABEL]]</div>
    <h1 class="hero-title">[[TITLE]]</h1>
    <p class="hero-sub">AI 분석 엔진이 원문 OT를 해체하여 시장·전략·크리에이티브·리스크의 4개 축으로 재구성한 기획 브리핑입니다.</p>
    <div class="hero-chips">
      <span class="chip chip-accent">📊 시장 분석</span>
      <span class="chip chip-accent">🎯 전략 논리</span>
      <span class="chip chip-accent">🎨 크리에이티브</span>
      <span class="chip chip-accent">⚠️ 리스크 점검</span>
      <span class="chip">🆔 [[JOB_ID]]</span>
      <span class="chip">📅 [[DATE]]</span>
    </div>
  </div>
  <div class="radar-wrap">
    <canvas id="radarChart" width="300" height="300"></canvas>
    <div class="radar-label">분석 커버리지 — 4개 축</div>
  </div>
</section>

<!-- Content sections -->
[[SECTIONS_HTML]]

<!-- Footer -->
<footer class="site-footer">
  <div class="footer-left">
    <div style="color:var(--text-2);font-weight:600;margin-bottom:.3rem">Antigravity Strategic Intelligence</div>
    <div>원본 파일: [[FILE_NAME]] · 생성일: [[DATE]]</div>
    <div style="margin-top:.5rem;color:var(--accent);font-size:.7rem">CONFIDENTIAL & PROPRIETARY</div>
  </div>
  <div class="footer-right">
    <span class="footer-chip">Claude Sonnet</span>
    <span class="footer-chip">v5.0</span>
  </div>
</footer>

<script>
// ── Scroll progress ────────────────────────
const bar = document.getElementById('scroll-progress');
window.addEventListener('scroll', () => {
  const pct = window.scrollY / (document.body.scrollHeight - window.innerHeight) * 100;
  bar.style.width = Math.min(pct, 100) + '%';
}, { passive: true });

// ── Intersection Observer ──────────────────
const io = new IntersectionObserver((entries) => {
  entries.forEach((e, idx) => {
    if (e.isIntersecting) {
      setTimeout(() => e.target.classList.add('visible'), idx * 80);
      io.unobserve(e.target);
    }
  });
}, { threshold: 0.08, rootMargin: '0px 0px -40px 0px' });
document.querySelectorAll('.ot-card').forEach(el => io.observe(el));

// ── Radar chart ────────────────────────────
const radarCtx = document.getElementById('radarChart').getContext('2d');
Chart.defaults.color = '#6e6e73';
Chart.defaults.font.family = "'Inter','Noto Sans KR',sans-serif";
new Chart(radarCtx, {
  type: 'radar',
  data: {
    labels: ['📊 시장 분석','🎯 전략 논리','🎨 크리에이티브','⚠️ 리스크'],
    datasets: [{
      label: '분석 완성도',
      data: [[RADAR_SCORES]],
      backgroundColor: 'rgba([[BRAND_RGB]],.15)',
      borderColor: 'rgba([[BRAND_RGB]],.8)',
      pointBackgroundColor: 'rgba([[BRAND_RGB]],1)',
      pointBorderColor: '#030303',
      pointBorderWidth: 2,
      pointRadius: 5,
      borderWidth: 2
    }]
  },
  options: {
    animation: { duration: 1200, easing: 'easeOutQuart' },
    scales: {
      r: {
        min: 0, max: 100,
        ticks: { display: false, stepSize: 25 },
        grid: { color: 'rgba(255,255,255,.06)' },
        angleLines: { color: 'rgba(255,255,255,.06)' },
        pointLabels: { font: { size: 11, weight: '500' }, color: '#a1a1a6', padding: 10 }
      }
    },
    plugins: {
      legend: { display: false },
      tooltip: {
        backgroundColor: '#111',
        borderColor: 'rgba([[BRAND_RGB]],.3)',
        borderWidth: 1,
        callbacks: { label: ctx => ` ${ctx.raw}점` }
      }
    }
  }
});

// ── Lucide icons ───────────────────────────
if (typeof lucide !== 'undefined') lucide.createIcons();
</script>
</body>
</html>"""

# ─── Main HTML generator ──────────────────────────────────────────────────────
def generate_html(data, output_dir, job_id, report_type="standard"):
    from datetime import datetime
    title = data.get("strategy_title", "OT 분석 리포트")
    categories = data.get("categories", {})
    file_name = data.get("file_name", "")
    date_str = datetime.now().strftime("%Y.%m.%d")

    _chart_counter[0] = 0

    sections_html = ""
    for cat_name, sections in categories.items():
        if not sections:
            continue
        sections_html += render_category(cat_name, sections)

    scores = compute_radar_scores(categories)
    theme = get_dynamic_theme(title)

    # Deep Intel: different badge label and output filename
    if report_type == "deep-intel":
        badge_label = "Deep Intelligence Report"
        out_filename = f"{job_id}_deep_intel_report.html"
    else:
        badge_label = "Strategic Intelligence Report"
        out_filename = f"{job_id}_precision_report.html"

    html = HTML_TEMPLATE
    html = html.replace("[[TITLE]]", title)
    html = html.replace("[[JOB_ID]]", job_id)
    html = html.replace("[[DATE]]", date_str)
    html = html.replace("[[FILE_NAME]]", file_name or "—")
    html = html.replace("[[SECTIONS_HTML]]", sections_html)
    html = html.replace("[[RADAR_SCORES]]", ", ".join(str(s) for s in scores))
    html = html.replace("[[BRAND_HEX]]", theme["hex"])
    html = html.replace("[[BRAND_RGB]]", theme["rgb"])
    html = html.replace("[[BADGE_LABEL]]", badge_label)

    os.makedirs(output_dir, exist_ok=True)
    out_path = os.path.join(output_dir, out_filename)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html)
    return out_path

# ─── Inline data builder (reads markdown files) ───────────────────────────────
def generate_strategic_base_inline(job_id, report_type="standard"):
    if report_type == "deep-intel":
        # Deep Intel: reads from 04_deep_intel directory
        source_path = ROOT / "outputs" / "04_deep_intel" / f"{job_id}_deep_intel_report.md"
        text = source_path.read_text(encoding="utf-8") if source_path.exists() else ""
    else:
        # Standard: reads from brain_engine deep analysis output
        try:
            from scripts.antigravity_bridge import load_markdown
            _, deep, _ = load_markdown(job_id)
        except:
            deep = ""
        text = deep or ""

    title = f"OT 분석 — {job_id}"
    # Find first H1 that isn't a job_id garbage header or OT-summary line
    for m in re.finditer(r'^#\s+(.+)', text, re.MULTILINE):
        candidate = m.group(1).strip()
        # Skip lines like "[20260414_xxx_ot] 초격차 딥-리포트" or "xxx — OT 분석 요약"
        if re.match(r'^\[.*_ot\]', candidate): continue
        if re.search(r'OT\s*분석\s*요약', candidate): continue
        title = candidate
        break

    file_name = ""
    mf = re.search(r'원본 파일[:\s]+(.+\.(?:pdf|docx|txt))', text, re.I)
    if mf:
        file_name = mf.group(1).strip()

    categories = {
        "Market Intelligence": [],
        "Strategic Logic": [],
        "Creative & Visual": [],
        "Gap Check & Risks": [],
    }

    sections = re.split(r'\n(?=##\s)', text)
    for sec in sections:
        lines = sec.strip().split('\n')
        if not lines:
            continue
        h2_m = re.match(r'^##\s+(.+)', lines[0])
        if not h2_m:
            continue
        h2 = h2_m.group(1).strip()

        cat = "Market Intelligence"
        if any(k in h2 for k in ["크리에이티브","Creative","Visual","비주얼"]): cat = "Creative & Visual"
        elif any(k in h2 for k in ["Gap","갭","리스크","Risk","질문","확인","불명확"]): cat = "Gap Check & Risks"
        elif any(k in h2 for k in ["전략","Strategic","Logic","Phase","단계"]): cat = "Strategic Logic"

        contents = [l.strip() for l in lines[1:] if l.strip()]
        if contents:
            categories[cat].append({
                "title": h2,
                "contents": contents,
                "importance": "high" if any(k in h2 for k in ["핵심","주요","Critical"]) else "normal"
            })

    return {"strategy_title": title, "categories": categories, "file_name": file_name}

# ─── PPTX (kept for compatibility, python-pptx optional) ─────────────────────
def generate_pptx(data, output_dir, job_id):
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print("[PPT Engine] python-pptx not installed. Skipping PPTX.")
        return None
    prs = Presentation()
    title_text = data.get("strategy_title", "OT Report")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    os.makedirs(output_dir, exist_ok=True)
    out = os.path.join(output_dir, f"{job_id}_strategic_deck.pptx")
    prs.save(out)
    return out

# ─── CLI entrypoint ───────────────────────────────────────────────────────────
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--job-id", required=True)
    parser.add_argument("-o", "--output", default="./outputs/03_final_reports")
    parser.add_argument("--report-type", default="standard",
                        choices=["standard", "deep-intel"],
                        help="standard: Claude 기반 OT 분석 | deep-intel: Gemini 기반 확장 리포트")
    args = parser.parse_args()

    try:
        data = generate_strategic_base_inline(args.job_id, report_type=args.report_type)
        label = "Deep Intel" if args.report_type == "deep-intel" else "Standard"
        print(f"[v5.0] Generating HTML ({label}) for {args.job_id}...")
        html_out = generate_html(data, args.output, args.job_id, report_type=args.report_type)
        print(f"DONE (HTML): {html_out}")

        # --- Cleanup intermediate markdown files ---
        try:
            summary_path = ROOT / "outputs" / "01_summaries" / f"{args.job_id}_summary.md"
            deep_path = ROOT / "outputs" / "02_deep_analysis" / f"{args.job_id}_deep_report.md"
            if summary_path.exists():
                summary_path.unlink()
                print(f"[Cleanup] Deleted: {summary_path.name}")
            if deep_path.exists():
                deep_path.unlink()
                print(f"[Cleanup] Deleted: {deep_path.name}")
        except Exception as cleanup_err:
            print(f"[Cleanup Error] {cleanup_err}")
    except Exception as e:
        import traceback; traceback.print_exc()
        sys.exit(1)
