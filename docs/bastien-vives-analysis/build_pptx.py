from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re, os

INPUT_MD = os.path.expanduser("~/research-output/바스티앙_비베스_작품분석/03_ppt_outline_revised.md")
OUTPUT_PPTX = os.path.expanduser("~/research-output/바스티앙_비베스_작품분석/04_presentation_draft.pptx")

COLOR_BG    = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TITLE = RGBColor(0x1a, 0x1a, 0x2e)
COLOR_BODY  = RGBColor(0x2d, 0x2d, 0x2d)
COLOR_ACCENT= RGBColor(0xe6, 0x39, 0x46)
COLOR_SLIDE_NUM = RGBColor(0x99, 0x99, 0x99)
COLOR_UNIT_BG   = RGBColor(0x1a, 0x1a, 0x2e)

W = Inches(13.333)
H = Inches(7.5)

def parse_outline(path):
    with open(path, encoding="utf-8") as f:
        content = f.read()
    # Stop at 수정 내역
    content = content.split("## 수정 내역")[0]

    slides = []
    current_unit = ""

    # Split by slide entries
    slide_pattern = re.compile(r'^\s*(\d+)\.\s+\*\*(.*?)\*\*\s*(?:\[(\w+)\])?', re.MULTILINE)

    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]
        unit_match = re.match(r'^## \[Unit (\d+)\] (.+)', line)
        if unit_match:
            current_unit = f"Unit {unit_match.group(1)}: {unit_match.group(2).split('(')[0].strip()}"
            i += 1
            continue

        slide_match = re.match(r'^\s*(\d+)\.\s+\*\*(.+?)\*\*', line)
        if slide_match:
            num = int(slide_match.group(1))
            raw_title = slide_match.group(2)
            layout = "Bento"
            layout_match = re.search(r'\[(\w+)\]', raw_title)
            if layout_match:
                layout = layout_match.group(1)
                raw_title = raw_title.replace(f'[{layout}]', '').strip()

            points = []
            i += 1
            while i < len(lines):
                pt_line = lines[i].strip()
                if pt_line.startswith('- '):
                    points.append(pt_line[2:].strip())
                    i += 1
                elif pt_line == '' or re.match(r'^\s*\d+\.', pt_line) or pt_line.startswith('##') or pt_line.startswith('*총'):
                    break
                else:
                    i += 1

            slides.append({
                'num': num,
                'title': raw_title,
                'points': points,
                'layout': layout,
                'unit': current_unit
            })
            continue
        i += 1

    return slides

def add_slide_number(slide, num, prs):
    txBox = slide.shapes.add_textbox(W - Inches(1), H - Inches(0.4), Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_SLIDE_NUM

def add_unit_tag(slide, unit_text):
    if not unit_text:
        return
    txBox = slide.shapes.add_textbox(Inches(0.4), H - Inches(0.45), Inches(5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = unit_text
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_SLIDE_NUM

def create_title_slide(prs, data):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_UNIT_BG

    # Accent line
    line = slide.shapes.add_shape(1, 0, Inches(3.2), W, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = data['title']
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Sub points
    y = Inches(3.8)
    for point in data['points'][:3]:
        txBox2 = slide.shapes.add_textbox(Inches(1.2), y, Inches(10), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = f"• {point}"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        y += Inches(0.45)

    add_slide_number(slide, data['num'], prs)
    return slide

def create_unit_cover(prs, data):
    """Unit 표지 슬라이드 (짝수 Unit 전환점에 삽입 안 함 — 일반 슬라이드로 처리)"""
    return create_standard_slide(prs, data)

def create_quote_slide(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)

    # Accent bar left
    bar = slide.shapes.add_shape(1, 0, Inches(1.5), Pt(6), Inches(4.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data['title']
    run.font.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = COLOR_TITLE

    # Quote (first point in big text)
    if data['points']:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(2.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = f'"{data["points"][0]}"'
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = COLOR_ACCENT

    # Remaining points
    y = Inches(4.2)
    for point in data['points'][1:]:
        txBox3 = slide.shapes.add_textbox(Inches(0.7), y, Inches(12), Inches(0.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = f"• {point}"
        run3.font.size = Pt(16)
        run3.font.color.rgb = COLOR_BODY
        y += Inches(0.45)

    add_slide_number(slide, data['num'], prs)
    add_unit_tag(slide, data['unit'])
    return slide

def create_table_slide(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_TITLE
    title_bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data['title']
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Points as table rows
    y = Inches(1.4)
    for idx, point in enumerate(data['points']):
        row_bg = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.55))
        if idx % 2 == 0:
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
        else:
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        row_bg.line.fill.background()

        txRow = slide.shapes.add_textbox(Inches(0.7), y + Pt(4), Inches(12), Inches(0.5))
        tf_row = txRow.text_frame
        p_row = tf_row.paragraphs[0]
        run_row = p_row.add_run()
        run_row.text = point
        run_row.font.size = Pt(15)
        run_row.font.color.rgb = COLOR_BODY
        y += Inches(0.6)
        if y > Inches(7.0):
            break

    add_slide_number(slide, data['num'], prs)
    add_unit_tag(slide, data['unit'])
    return slide

def create_timeline_slide(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data['title']
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = COLOR_TITLE

    # Vertical timeline line
    line = slide.shapes.add_shape(1, Inches(1.5), Inches(1.2), Pt(3), Inches(5.8))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Timeline items
    y = Inches(1.3)
    step = Inches(5.5) / max(len(data['points']), 1)
    for idx, point in enumerate(data['points']):
        # Dot
        dot = slide.shapes.add_shape(9, Inches(1.35), y + Inches(0.05), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_ACCENT
        dot.line.fill.background()

        # Text
        txItem = slide.shapes.add_textbox(Inches(2.0), y, Inches(10.5), Inches(0.6))
        tf_item = txItem.text_frame
        tf_item.word_wrap = True
        p_item = tf_item.paragraphs[0]
        run_item = p_item.add_run()
        run_item.text = point
        run_item.font.size = Pt(16)
        run_item.font.color.rgb = COLOR_BODY
        y += step

    add_slide_number(slide, data['num'], prs)
    add_unit_tag(slide, data['unit'])
    return slide

def create_standard_slide(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)
    title_bar.line.fill.background()

    # Accent line under title
    accent_line = slide.shapes.add_shape(1, 0, Inches(1.2), Inches(2), Pt(3))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLOR_ACCENT
    accent_line.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data['title']
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = COLOR_TITLE

    # Content points — bento grid style (2 columns if 4+ points)
    points = data['points']
    if len(points) >= 4:
        col1 = points[:len(points)//2 + len(points)%2]
        col2 = points[len(points)//2 + len(points)%2:]

        for col_idx, col_pts in enumerate([col1, col2]):
            x = Inches(0.5) + col_idx * Inches(6.4)
            y = Inches(1.5)
            for point in col_pts:
                card = slide.shapes.add_shape(1, x, y, Inches(6.0), Inches(0.7))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFF)
                card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xFF)

                txPt = slide.shapes.add_textbox(x + Inches(0.15), y + Pt(6), Inches(5.7), Inches(0.6))
                tf_pt = txPt.text_frame
                tf_pt.word_wrap = True
                p_pt = tf_pt.paragraphs[0]
                run_pt = p_pt.add_run()
                run_pt.text = f"▸  {point}"
                run_pt.font.size = Pt(15)
                run_pt.font.color.rgb = COLOR_BODY
                y += Inches(0.8)
    else:
        y = Inches(1.6)
        for point in points:
            card = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.8))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFF)
            card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xFF)

            txPt = slide.shapes.add_textbox(Inches(0.7), y + Pt(8), Inches(12.0), Inches(0.65))
            tf_pt = txPt.text_frame
            tf_pt.word_wrap = True
            p_pt = tf_pt.paragraphs[0]
            run_pt = p_pt.add_run()
            run_pt.text = f"▸  {point}"
            run_pt.font.size = Pt(16)
            run_pt.font.color.rgb = COLOR_BODY
            y += Inches(0.9)

    add_slide_number(slide, data['num'], prs)
    add_unit_tag(slide, data['unit'])
    return slide

def build(input_md, output_pptx):
    slides_data = parse_outline(input_md)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for data in slides_data:
        layout = data.get('layout', 'Bento')
        num = data['num']

        if num == 1:
            create_title_slide(prs, data)
        elif layout == 'Quote':
            create_quote_slide(prs, data)
        elif layout == 'Table':
            create_table_slide(prs, data)
        elif layout == 'Timeline':
            create_timeline_slide(prs, data)
        else:
            create_standard_slide(prs, data)

    prs.save(output_pptx)
    print(f"PPTX 생성 완료: {output_pptx} ({len(slides_data)}슬라이드)")

if __name__ == "__main__":
    build(INPUT_MD, OUTPUT_PPTX)
